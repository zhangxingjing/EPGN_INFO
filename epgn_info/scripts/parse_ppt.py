import re
import base64
import time
from pprint import pprint

from pptx import Presentation
from pptx.util import Inches, Pt
from epgn_info.epgn_info.settings.dev import BASE_DIR    # Nginx
# from epgn_info.settings.devp import BASE_DIR  # manage

PPT_MODEL_PATH = BASE_DIR + '/apps/calculate/PPTModel/'
IMAGE_LOCATION = [
    ['F3 VZ', '(N)G3 VZ'],
    ['F3 VS', '(N)G3 VS'],
    ['F5 VZ', '(N)G5 VZ'],
    ['KP 80-20'],
    ['(Square&Lab)Leerlauf D Gang mit AC'],
    ['(Square&Lab)Leerlauf D Gang ohne Verbrauche'],
    ['(Square&Lab)Leerlauf D Gang mit Verbrauche'],
    ['(Square&Lab)Leerlauf R Gang mit AC'],
]


class ParsePPT():
    def __init__(self, item):
        """
        处理前端的数据
        :param item:前端传来的图片的位置信息
        """
        self.filename = time.strftime('%Y_%m_%d_%H_%M_%S', time.localtime(time.time()))
        self.prs = Presentation(PPT_MODEL_PATH + "4zuo.pptx")
        # self.save_path = save_path
        self.data = item
        self.img_list = []

    def parse_title(self):
        """
        处理当前PPT首页的标题
        :return: 直接操作全局变量，不返回
        """
        # 拿到第一页幻灯片
        # 把前端获取到的数据，填入幻灯片中

        blank_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        # 设置要新建的文本框的位置
        left = top = width = height = Inches(1)
        # 实例化一个文本框
        txBox = slide.shapes.add_textbox(left, top, width, height)
        # 设置文件框的类型
        tf = txBox.text_frame
        # 给定文本框里的文字
        # tf.text = self.title
        # 添加段落，向下在添加段落文字
        p = tf.add_paragraph()
        # 给新增加的段落添加文字
        p.text = "This is a second add_paragraph that's bold"
        # 给新添加的段落文字设置为粗体
        p.font.bold = True
        # 再在这个文本框中新建一个段落
        p = tf.add_paragraph()
        # 设置新段落的文字
        p.text = "This is a third paragraph that's big"
        # 设置新添加的段落文字的字号为40
        p.font.size = Pt(40)

    def parse_pic(self):
        """
        处理前端返回的算法结果的图片 ==> base64--pic
        :return: 返回图片存储路径（图片使用完成后是否删除当前生成的图片）
        """
        """
        for noise in self.data["internal_noise"]:  # 内部噪声
            pic_info = noise["info"]
            for key, value in noise["data"].items():
                str_bas64 = re.match(r'(data:image/png;base64,(.*))', value).group(2)
                img = base64.b64decode(str_bas64)
                pic_path = self.save_path + 'image/{}.jpg'.format(pic_info + " " + key)
                fh = open(PPT_MODEL_PATH + pic_path, 'wb')
                fh.write(img)
                fh.close()
                if key == "VL":
                    self.vl[pic_info] = pic_path
                elif key == "VR":
                    self.vr[pic_info] = pic_path
                elif key == "HL":
                    self.hl[pic_info] = pic_path
                elif key == "HR":
                    self.hr[pic_info] = pic_path
                elif key == "ML":
                    self.ml[pic_info] = pic_path
                elif key == "MR":
                    self.mr[pic_info] = pic_path

        pic_gun_info = self.data["roolgeraeush_noise"]["info"]
        for key, value in self.data["roolgeraeush_noise"]["data"].items():
            str_bas64 = re.match(r'(data:image/png;base64,(.*))', value).group(2)
            img = base64.b64decode(str_bas64)
            pic_path = self.save_path + 'image/{}.jpg'.format(pic_gun_info + " " + key)
            fh = open(PPT_MODEL_PATH + pic_path, 'wb')
            fh.write(img)
            fh.close()
            if key == "VL":
                self.vl[pic_gun_info] = pic_path
            elif key == "VR":
                self.vr[pic_gun_info] = pic_path
            elif key == "HL":
                self.hl[pic_gun_info] = pic_path
            elif key == "HR":
                self.hr[pic_gun_info] = pic_path
        """

        # 在这里处理当前PPT，生成之后返回当前的PPT路径
        for item in self.data["items"]:
            pic_info = re.search(r' (.*)', item["status"]).group(1) + "_" + item["filename"] # 正则到合适的imgName
            str_bas64 = re.match(r'(data:image/png;base64,(.*))', item["base64"]).group(2)
            img = base64.b64decode(str_bas64)
            pic_path = PPT_MODEL_PATH + 'image/{}.jpg'.format(pic_info)
            fh = open(pic_path, 'wb')
            fh.write(img)
            # fh.close()

            # yield item["status"], item["channel"], pic_path
            img_dict = {
                "channel": item["filename"],
                "status": item["status"],
                "pic_path": pic_path
            }
            yield img_dict
            # self.img_list.append(img_dict)

    def insert_pic(self, img_dict):
        """
        通过处理后的前端数据，合成规定的报告
        :return:处理完成的PPT
        """

        # TODO: 后面多个工况在一起的情况怎么处理的？

        """
        if self.ml:  # 要么是4坐， 要么就是6坐 ==> 复制幻灯片
            self.prs = Presentation(PPT_MODEL_PATH + "PPTModel/6zuo.pptx")
            for key in key_list:
                num = key_list.index(key)
                # 左前
                left, top, width, height = Inches(0.5), Inches(1.6), Inches(4.5), Inches(2.4)
                self.prs.slides[2 * num + 1].shapes.add_picture('{}'.format(PPT_MODEL_PATH + self.vl[key]), left, top,
                                                                width, height)
                # 右前
                left, top, width, height = Inches(5), Inches(1.6), Inches(4.5), Inches(2.4)
                self.prs.slides[2 * num + 1].shapes.add_picture('{}'.format(PPT_MODEL_PATH + self.vr[key]), left, top,
                                                                width, height)
                # 左中
                left, top, width, height = Inches(0.5), Inches(4), Inches(4.5), Inches(2.4)
                self.prs.slides[2 * num + 1].shapes.add_picture('{}'.format(PPT_MODEL_PATH + self.hl[key]), left, top,
                                                                width, height)
                # 右中
                left, top, width, height = Inches(5), Inches(4), Inches(4.5), Inches(2.4)
                self.prs.slides[2 * num + 1].shapes.add_picture('{}'.format(PPT_MODEL_PATH + self.hr[key]), left, top,
                                                                width, height)
                # 左后
                left, top, width, height = Inches(0.5), Inches(1.6), Inches(4.5), Inches(2.4)
                self.prs.slides[2 * num + 2].shapes.add_picture('{}'.format(PPT_MODEL_PATH + self.vl[key]), left, top,
                                                                width, height)
                # 右后
                left, top, width, height = Inches(5), Inches(1.6), Inches(4.5), Inches(2.4)
                self.prs.slides[2 * num + 2].shapes.add_picture('{}'.format(PPT_MODEL_PATH + self.vr[key]), left, top,
                                                                width, height)
        else:
            for key in key_list:
                num = img_dict.index(key)
                # 左前
                left, top, width, height = Inches(0.5), Inches(1.6), Inches(4.5), Inches(2.4)
                self.prs.slides[num + 1].shapes.add_picture('{}'.format(img_dict[key]), left, top, width,height)
                # 右前
                left, top, width, height = Inches(5), Inches(1.6), Inches(4.5), Inches(2.4)
                self.prs.slides[num + 1].shapes.add_picture('{}'.format(img_dict[key]), left, top, width,height)
                # 左后
                left, top, width, height = Inches(0.5), Inches(4), Inches(4.5), Inches(2.4)
                self.prs.slides[num + 1].shapes.add_picture('{}'.format(img_dict[key]), left, top, width,height)
                # 右后
                left, top, width, height = Inches(5), Inches(4), Inches(4.5), Inches(2.4)
                self.prs.slides[num + 1].shapes.add_picture('{}'.format(img_dict[key]), left, top, width, height)
    """

        print(img_dict["channel"], img_dict["status"])

        # 获取当前图片在幻灯片中的位置，将其放置在对应位置上
        num = 0
        for key in IMAGE_LOCATION:  # 使用全局变量
            if img_dict["status"] in key:
                num = IMAGE_LOCATION.index(key) + 2 -1

        if img_dict["channel"] == "VL" or img_dict["channel"] == "vorn links":
            # 左前
            left, top, width, height = Inches(0.5), Inches(1.6), Inches(4.5), Inches(2.4)
            self.prs.slides[num].shapes.add_picture('{}'.format(img_dict["pic_path"]), left, top, width, height)
        elif img_dict["channel"] == "VR" or img_dict["channel"] == "vorn rechits" or img_dict["channel"] == "LR_X":
            # 右前
            left, top, width, height = Inches(5), Inches(1.6), Inches(4.5), Inches(2.4)
            self.prs.slides[num].shapes.add_picture('{}'.format(img_dict["pic_path"]), left, top, width, height)
        elif img_dict["channel"] == "HL" or img_dict["channel"] == "hinten links" or img_dict["channel"] == "LR_Y":
            # 左后
            left, top, width, height = Inches(0.5), Inches(4), Inches(4.5), Inches(2.4)
            self.prs.slides[num].shapes.add_picture('{}'.format(img_dict["pic_path"]), left, top, width, height)
        elif img_dict["channel"] == "HR" or img_dict["channel"] == "hinten rechits" or img_dict["channel"] == "LR_Z":
            # 右后
            left, top, width, height = Inches(5), Inches(4), Inches(4.5), Inches(2.4)
            self.prs.slides[num].shapes.add_picture('{}'.format(img_dict["pic_path"]), left, top, width, height)
        elif img_dict["channel"] == "FS_IN_X":
            left, top, width, height = Inches(0.5), Inches(1.6), Inches(4.5), Inches(2.4)
            self.prs.slides[num + 1].shapes.add_picture('{}'.format(img_dict["pic_path"]), left, top, width, height)
        elif img_dict["channel"] == "FS_IN_Z":
            left, top, width, height = Inches(5), Inches(1.6), Inches(4.5), Inches(2.4)
            self.prs.slides[num + 1].shapes.add_picture('{}'.format(img_dict["pic_path"]), left, top, width, height)

    def save_ppt(self):
        """
        保存最终的PPT文件
        :return:PPT存储位置
        """
        self.prs.save(PPT_MODEL_PATH + 'ppt_result/{}.pptx'.format(self.filename))
        return PPT_MODEL_PATH + 'ppt_result/{}.pptx'.format(self.filename)

    def run(self):
        """
        面向对象的接口
        :return: 当前PPT存储的位置
        """
        # try:
            # self.parse_title()
        for img_dict in self.parse_pic():
            pprint(img_dict)
            self.insert_pic(img_dict)
        ppt_path = self.save_ppt()
        return ppt_path
        # except Exception as e:
        #     print(e)
        #     raise Exception
